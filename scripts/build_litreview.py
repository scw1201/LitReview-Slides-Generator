#!/usr/bin/env python3
import argparse
import json
import os
import re
import shutil
import socket
import sqlite3
import subprocess
import sys
import tempfile
import time
import urllib.error
import urllib.request
from collections import Counter, defaultdict
from dataclasses import dataclass
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple

try:
    import fitz  # PyMuPDF
except Exception:
    fitz = None

try:
    from sklearn.cluster import KMeans
    from sklearn.feature_extraction.text import TfidfVectorizer
    from sklearn.metrics import silhouette_score
except Exception:
    KMeans = None
    TfidfVectorizer = None
    silhouette_score = None

try:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.util import Inches, Pt
except Exception:
    Presentation = None


SEED = 42

DEFAULT_SECTION_PATTERNS: Dict[str, List[str]] = {
    "abstract": [r"^abstract$"],
    "intro": [r"\bintroduction\b", r"\bbackground\b", r"\bmotivation\b"],
    "method": [
        r"\bmaterials?\s+and\s+methods?\b",
        r"\bmethods?\b",
        r"\bmethodology\b",
        r"\bapproach\b",
        r"\bframework\b",
        r"\bmodel\b",
        r"\bimplementation\b",
    ],
    "results": [r"\bexperiment", r"\bresults?\b", r"\bevaluation\b", r"\buser study\b", r"\bablation\b"],
    "discussion": [r"\bdiscussion\b"],
    "conclusion": [r"\bconclusion\b", r"\bconcluding\b", r"\bconclusion\s+and\s+future\s+work\b"],
    "limitations": [r"\blimitation", r"\bfuture work\b", r"\bfuture directions?\b", r"\bthreats to validity\b"],
}

SECTION_PATTERNS: Dict[str, List[str]] = dict(DEFAULT_SECTION_PATTERNS)


@dataclass
class SentenceSpan:
    text: str
    page: int


class LLMClient:
    def __init__(
        self,
        mode: str,
        model: str,
        base_url: str,
        api_key: Optional[str],
        codex_bin: Optional[str],
        timeout_sec: int,
        max_input_chars: int,
        max_tokens: int,
    ):
        self.mode = mode
        self.model = model
        self.base_url = base_url.rstrip("/")
        self.api_key = api_key
        self.codex_bin = self._resolve_codex_bin(codex_bin)
        self.timeout_sec = timeout_sec
        self.max_input_chars = max_input_chars
        self.max_tokens = max_tokens
        self.last_error = ""
        self.last_raw_content = ""

    @property
    def enabled(self) -> bool:
        if self.mode == "codex_cli":
            return bool(self.codex_bin)
        return self.mode == "openai_compatible" and bool(self.api_key)

    def _resolve_codex_bin(self, configured: Optional[str]) -> str:
        c = (configured or "").strip()
        if c and Path(c).exists():
            return c
        found = shutil.which("codex")
        if found:
            return found
        candidates = [
            "/Applications/Codex.app/Contents/Resources/codex",
            str(Path.home() / ".local" / "bin" / "codex"),
        ]
        for x in candidates:
            if Path(x).exists():
                return x
        return "codex"

    def _extract_json_obj(self, text: str) -> Optional[Dict[str, Any]]:
        text = (text or "").strip()
        if not text:
            return None
        if text.startswith("```"):
            text = re.sub(r"^```(?:json)?\s*", "", text)
            text = re.sub(r"\s*```$", "", text)
        try:
            return json.loads(text)
        except Exception:
            pass
        decoder = json.JSONDecoder()
        for i, ch in enumerate(text):
            if ch != "{":
                continue
            try:
                obj, _ = decoder.raw_decode(text[i:])
                if isinstance(obj, dict):
                    return obj
            except Exception:
                continue
        m = re.search(r"\{.*\}", text, re.S)
        if not m:
            return None
        try:
            return json.loads(m.group(0))
        except Exception:
            return None

    def _run_codex_exec(self, prompt: str) -> Optional[str]:
        with tempfile.TemporaryDirectory(prefix="litreview_codex_") as td:
            out_file = Path(td) / "last_message.txt"
            cmd = [
                self.codex_bin,
                "exec",
                prompt[: self.max_input_chars],
                "--model",
                self.model,
                "--skip-git-repo-check",
                "--full-auto",
                "--ephemeral",
                "--output-last-message",
                str(out_file),
                "--color",
                "never",
            ]
            try:
                proc = subprocess.run(
                    cmd,
                    capture_output=True,
                    text=True,
                    timeout=self.timeout_sec + 30,
                    check=False,
                )
            except subprocess.TimeoutExpired:
                self.last_error = "TimeoutError: codex_cli timeout"
                return None
            except Exception as e:
                self.last_error = f"{type(e).__name__}: {e}"
                return None

            content = ""
            if out_file.exists():
                try:
                    content = out_file.read_text(encoding="utf-8").strip()
                except Exception:
                    content = ""
            if not content:
                content = (proc.stdout or "").strip()

            self.last_raw_content = content
            if proc.returncode != 0:
                err = (proc.stderr or "").strip() or f"codex_cli exit={proc.returncode}"
                if content:
                    self.last_error = f"codex_cli exit={proc.returncode}"
                else:
                    self.last_error = err
            return content or None

    def _post_chat(self, payload: Dict[str, Any]) -> Dict[str, Any]:
        data = json.dumps(payload).encode("utf-8")
        req = urllib.request.Request(
            f"{self.base_url}/chat/completions",
            data=data,
            method="POST",
            headers={
                "Content-Type": "application/json",
                "Authorization": f"Bearer {self.api_key}",
            },
        )
        with urllib.request.urlopen(req, timeout=self.timeout_sec) as resp:
            return json.loads(resp.read().decode("utf-8"))

    def complete_json(self, system_prompt: str, user_prompt: str) -> Optional[Dict[str, Any]]:
        if not self.enabled:
            return None
        self.last_error = ""
        self.last_raw_content = ""
        if self.mode == "codex_cli":
            prompt = (
                "你是严格JSON输出器。仅输出一个JSON对象，不要输出额外解释、代码块、标题。\n"
                f"[SYSTEM]\n{system_prompt}\n\n[USER]\n{user_prompt}"
            )
            content = self._run_codex_exec(prompt)
            if not content:
                if not self.last_error:
                    self.last_error = "empty_content"
                return None
            parsed = self._extract_json_obj(content)
            if parsed is not None:
                return parsed
            self.last_error = "json_parse_failed"
            return None
        messages = [
            {"role": "system", "content": system_prompt},
            {"role": "user", "content": user_prompt[: self.max_input_chars]},
        ]
        # Try strict JSON response first.
        payloads = [
            {
                "model": self.model,
                "temperature": 0.2,
                "max_tokens": self.max_tokens,
                "response_format": {"type": "json_object"},
                "messages": messages,
            },
            {
                "model": self.model,
                "temperature": 0.2,
                "max_tokens": self.max_tokens,
                "messages": messages,
            },
        ]
        try:
            for payload in payloads:
                try:
                    body = self._post_chat(payload)
                    content = (
                        body.get("choices", [{}])[0]
                        .get("message", {})
                        .get("content", "")
                        .strip()
                    )
                    self.last_raw_content = content
                    parsed = self._extract_json_obj(content)
                    if parsed is not None:
                        return parsed
                    self.last_error = "json_parse_failed"
                except (
                    urllib.error.URLError,
                    urllib.error.HTTPError,
                    json.JSONDecodeError,
                    KeyError,
                    ValueError,
                    TimeoutError,
                    socket.timeout,
                ) as e:
                    self.last_error = f"{type(e).__name__}: {e}"
                    continue
            return None
        except Exception as e:
            self.last_error = f"{type(e).__name__}: {e}"
            return None

    def complete_text(self, system_prompt: str, user_prompt: str) -> Optional[str]:
        if not self.enabled:
            return None
        self.last_error = ""
        self.last_raw_content = ""
        if self.mode == "codex_cli":
            prompt = (
                "请直接输出最终文本，不要附加解释。\n"
                f"[SYSTEM]\n{system_prompt}\n\n[USER]\n{user_prompt}"
            )
            content = self._run_codex_exec(prompt)
            if not content:
                if not self.last_error:
                    self.last_error = "empty_content"
                return None
            return content
        payload = {
            "model": self.model,
            "temperature": 0.2,
            "max_tokens": self.max_tokens,
            "messages": [
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": user_prompt[: self.max_input_chars]},
            ],
        }
        try:
            body = self._post_chat(payload)
            content = (
                body.get("choices", [{}])[0]
                .get("message", {})
                .get("content", "")
                .strip()
            )
            self.last_raw_content = content
            if not content:
                self.last_error = "empty_content"
                return None
            return content
        except (
            urllib.error.URLError,
            urllib.error.HTTPError,
            json.JSONDecodeError,
            KeyError,
            ValueError,
            TimeoutError,
            socket.timeout,
        ) as e:
            self.last_error = f"{type(e).__name__}: {e}"
        return None


class ZoteroMcpRagClient:
    def __init__(
        self,
        enabled: bool,
        python_bin: str,
        config_path: str,
        use_local: bool = True,
        timeout_sec: int = 120,
    ):
        self.enabled = enabled
        self.python_bin = self._resolve_python_bin(python_bin)
        self.config_path = config_path.strip()
        self.use_local = use_local
        self.timeout_sec = timeout_sec
        self.last_error = ""

    def _resolve_python_bin(self, configured: str) -> str:
        c = (configured or "").strip()
        if c and Path(c).exists():
            return c
        # Preferred known installation path from uv tool.
        preferred = str(Path.home() / ".local" / "share" / "uv" / "tools" / "zotero-mcp" / "bin" / "python")
        if Path(preferred).exists():
            return preferred
        return sys.executable

    def search(self, query: str, limit: int) -> List[Dict[str, Any]]:
        self.last_error = ""
        if not self.enabled:
            return []
        q = (query or "").strip()
        if not q:
            return []
        script = r"""
import json, sys
from zotero_mcp.semantic_search import ZoteroSemanticSearch
query = sys.argv[1]
limit = int(sys.argv[2])
config_path = sys.argv[3] if len(sys.argv) > 3 else ""
if config_path:
    s = ZoteroSemanticSearch(config_path=config_path)
else:
    s = ZoteroSemanticSearch()
r = s.search(query, limit=limit)
print(json.dumps(r, ensure_ascii=False))
"""
        cmd = [self.python_bin, "-c", script, q, str(max(1, int(limit))), self.config_path]
        run_env = os.environ.copy()
        if self.use_local:
            run_env["ZOTERO_LOCAL"] = "true"
            run_env.setdefault("NO_PROXY", "localhost,127.0.0.1,::1")
            run_env.setdefault("no_proxy", "localhost,127.0.0.1,::1")
        try:
            proc = subprocess.run(
                cmd,
                capture_output=True,
                text=True,
                timeout=self.timeout_sec,
                check=False,
                env=run_env,
            )
        except subprocess.TimeoutExpired:
            self.last_error = "rag_timeout"
            return []
        except Exception as e:
            self.last_error = f"{type(e).__name__}: {e}"
            return []

        if proc.returncode != 0:
            err = (proc.stderr or "").strip()
            self.last_error = err or f"rag_exit_{proc.returncode}"
            return []
        try:
            payload = json.loads((proc.stdout or "").strip() or "{}")
        except Exception as e:
            self.last_error = f"rag_json_parse_error: {e}"
            return []
        return self._normalize_results(payload)

    def _normalize_results(self, payload: Dict[str, Any]) -> List[Dict[str, Any]]:
        results = payload.get("results") if isinstance(payload, dict) else None
        if not isinstance(results, list):
            return []
        out: List[Dict[str, Any]] = []
        for r in results:
            if not isinstance(r, dict):
                continue
            meta = r.get("metadata", {}) if isinstance(r.get("metadata"), dict) else {}
            title = str(r.get("title") or meta.get("title") or "").strip()
            snippet = str(
                r.get("snippet")
                or r.get("content")
                or r.get("document")
                or meta.get("snippet")
                or ""
            ).strip()
            item_key = str(r.get("item_key") or meta.get("item_key") or meta.get("key") or "").strip()
            score_raw = r.get("score", r.get("similarity", r.get("distance")))
            score = None
            try:
                score = float(score_raw)
            except Exception:
                score = None
            if not title and not snippet:
                continue
            out.append(
                {
                    "title": title,
                    "snippet": snippet,
                    "item_key": item_key,
                    "score": score,
                }
            )
        return out


def parse_args() -> argparse.Namespace:
    p = argparse.ArgumentParser(
        description="Generate literature-review PPT/MD/JSON from local Zotero PDF manifest."
    )
    p.add_argument("--collection", required=True, help="Collection name or key")
    p.add_argument("--mode", choices=["analyze", "global", "render", "all"], default="all")
    p.add_argument("--manifest", help="Optional path to Zotero item manifest JSON")
    p.add_argument("--max_papers", type=int, default=20)
    p.add_argument("--cluster_k", type=int, default=0, help="0 means auto clustering")
    p.add_argument("--language", choices=["zh", "en"], default="zh")
    p.add_argument("--include_images", type=str, default="true")
    p.add_argument("--llm_mode", choices=["off", "openai_compatible", "codex_cli"], default="codex_cli")
    p.add_argument("--llm_model", default="gpt-5-mini")
    p.add_argument("--llm_base_url", default="https://api.openai.com/v1")
    p.add_argument("--llm_api_key_env", default="OPENAI_API_KEY")
    p.add_argument("--codex_bin", default="", help="Path to codex executable for llm_mode=codex_cli")
    p.add_argument("--llm_timeout_sec", type=int, default=180)
    p.add_argument("--llm_max_input_chars", type=int, default=12000)
    p.add_argument("--llm_max_tokens", type=int, default=180)
    p.add_argument("--rag_enabled", type=str, default="false")
    p.add_argument("--rag_top_k", type=int, default=8)
    p.add_argument("--rag_python_bin", default="", help="Python path with zotero_mcp installed")
    p.add_argument("--rag_config_path", default="", help="Optional semantic search config path for zotero-mcp")
    p.add_argument("--rag_use_local", type=str, default="true", help="Set ZOTERO_LOCAL=true for RAG subprocess")
    p.add_argument("--zotero_db", default=str(Path.home() / "Zotero" / "zotero.sqlite"))
    p.add_argument("--zotero_storage_dir", default=str(Path.home() / "Zotero" / "storage"))
    p.add_argument("--output_dir", default="outputs")
    p.add_argument("--session_name", default="", help="Session id/name for this run")
    p.add_argument(
        "--session_layout",
        choices=["folder", "filename"],
        default="folder",
        help="folder: output_dir/session_name, filename: append session_name to files",
    )
    p.add_argument("--analyze_json", default="", help="Path to per-paper analysis JSON")
    p.add_argument("--global_json", default="", help="Path to global synthesis JSON")
    p.add_argument("--config_json", default="", help="Optional config JSON; CLI args override missing keys only")
    p.add_argument("--status_file", default="")
    p.add_argument("--log_file", default="")
    p.add_argument("--paper_status_file", default="")
    p.add_argument("--section_map_json", default="", help="Optional JSON file to override section regex map")
    p.add_argument("--verbose", action="store_true")
    return p.parse_args()


def parse_bool(value: str) -> bool:
    return str(value).strip().lower() in {"1", "true", "yes", "y", "on"}


def ts_now() -> str:
    return time.strftime("%Y-%m-%d %H:%M:%S")


def append_log(log_file: Optional[Path], msg: str) -> None:
    if log_file is None:
        return
    line = f"[{ts_now()}] {msg}\n"
    try:
        with log_file.open("a", encoding="utf-8") as f:
            f.write(line)
    except Exception:
        pass


def write_status(
    status_file: Optional[Path],
    *,
    stage: str,
    status: str,
    progress_current: int,
    progress_total: int,
    message: str,
    extra: Optional[Dict[str, Any]] = None,
) -> None:
    if status_file is None:
        return
    data: Dict[str, Any] = {
        "timestamp": ts_now(),
        "stage": stage,
        "status": status,
        "progress_current": progress_current,
        "progress_total": progress_total,
        "progress_pct": 0 if progress_total <= 0 else round(progress_current * 100.0 / progress_total, 1),
        "message": message,
    }
    if extra:
        data.update(extra)
    tmp = status_file.with_suffix(status_file.suffix + ".tmp")
    try:
        tmp.write_text(json.dumps(data, ensure_ascii=False, indent=2), encoding="utf-8")
        tmp.replace(status_file)
    except Exception:
        pass


def append_jsonl(path: Optional[Path], obj: Dict[str, Any]) -> None:
    if path is None:
        return
    try:
        with path.open("a", encoding="utf-8") as f:
            f.write(json.dumps(obj, ensure_ascii=False) + "\n")
    except Exception:
        pass


def load_config_json(path: Optional[Path]) -> Dict[str, Any]:
    if path is None or not str(path).strip():
        return {}
    if not path.exists():
        raise ValueError(f"config json not found: {path}")
    data = json.loads(path.read_text(encoding="utf-8"))
    if not isinstance(data, dict):
        raise ValueError("config json must be an object")
    return data


def apply_config_defaults(args: argparse.Namespace, cfg: Dict[str, Any]) -> None:
    # Fill only when CLI value equals parser default.
    defaults = {
        "mode": "all",
        "max_papers": 20,
        "cluster_k": 0,
        "language": "zh",
        "include_images": "true",
        "llm_mode": "codex_cli",
        "llm_model": "gpt-5-mini",
        "llm_base_url": "https://api.openai.com/v1",
        "codex_bin": "",
        "llm_timeout_sec": 180,
        "llm_max_input_chars": 12000,
        "llm_max_tokens": 180,
        "rag_enabled": "false",
        "rag_top_k": 8,
        "rag_python_bin": "",
        "rag_config_path": "",
        "rag_use_local": "true",
        "output_dir": "outputs",
        "session_name": "",
        "session_layout": "folder",
        "section_map_json": "",
    }
    for k, default_v in defaults.items():
        if k in cfg and hasattr(args, k) and getattr(args, k) == default_v:
            setattr(args, k, cfg[k])


def validate_args(args: argparse.Namespace) -> None:
    if args.max_papers < 1 or args.max_papers > 200:
        raise ValueError("--max_papers must be in [1, 200]")
    if args.manifest:
        manifest = Path(args.manifest)
        if not manifest.exists():
            raise ValueError(f"manifest not found: {manifest}")
    if args.llm_timeout_sec < 5:
        raise ValueError("--llm_timeout_sec must be >= 5")
    if args.llm_max_input_chars < 1000:
        raise ValueError("--llm_max_input_chars must be >= 1000")
    if args.llm_max_tokens < 32:
        raise ValueError("--llm_max_tokens must be >= 32")
    if args.rag_top_k < 1 or args.rag_top_k > 50:
        raise ValueError("--rag_top_k must be in [1, 50]")
    if args.cluster_k < 0:
        raise ValueError("--cluster_k must be >= 0")
    if args.section_map_json:
        p = Path(args.section_map_json)
        if not p.exists():
            raise ValueError(f"section map not found: {p}")


def sanitize_session_name(name: str) -> str:
    text = (name or "").strip()
    if not text:
        return ""
    return re.sub(r"[^A-Za-z0-9._-]+", "-", text).strip("._-")


def resolve_session_context(args: argparse.Namespace) -> Tuple[Path, str, str]:
    collection_dir = sanitize_session_name(str(getattr(args, "collection", ""))) or str(getattr(args, "collection", "collection"))
    session_name = sanitize_session_name(getattr(args, "session_name", ""))
    base_stem = f"review_{args.collection}"
    output_root = Path(args.output_dir)
    if str(getattr(args, "session_layout", "folder")) == "filename":
        if not session_name:
            session_name = collection_dir
        return output_root, f"{base_stem}.{session_name}", session_name
    # Default folder layout: one collection => one folder.
    # output_dir/<collection>/...
    if not session_name or session_name == collection_dir:
        return output_root / collection_dir, base_stem, collection_dir
    # Optional sub-session under collection folder.
    # output_dir/<collection>/<session>/...
    return output_root / collection_dir / session_name, base_stem, session_name


def load_section_patterns(path: Optional[Path]) -> Dict[str, List[str]]:
    if path is None or not str(path).strip():
        return dict(DEFAULT_SECTION_PATTERNS)
    data = json.loads(path.read_text(encoding="utf-8"))
    if not isinstance(data, dict):
        raise ValueError("section map must be a JSON object")
    out: Dict[str, List[str]] = {}
    for k, v in data.items():
        if not isinstance(k, str) or not isinstance(v, list):
            continue
        pats = [str(x).strip() for x in v if str(x).strip()]
        if pats:
            out[k] = pats
    if not out:
        raise ValueError("section map has no valid patterns")
    return out


def read_manifest(path: Path) -> Dict:
    with path.open("r", encoding="utf-8") as f:
        data = json.load(f)
    if not isinstance(data, dict):
        raise ValueError("manifest must be a JSON object")
    if "items" not in data or not isinstance(data["items"], list):
        raise ValueError("manifest.items must be a list")
    if "collection" not in data:
        data["collection"] = "unknown"
    return data


def choose_first_nonempty(values: List[Optional[str]], default: str = "unknown") -> str:
    for v in values:
        if v is None:
            continue
        text = str(v).strip()
        if text:
            return text
    return default


def parse_year(date_text: str) -> str:
    if not date_text:
        return "unknown"
    m = re.search(r"(19|20)\d{2}", str(date_text))
    if m:
        return m.group(0)
    return "unknown"


def resolve_pdf_path(raw_path: str, attachment_key: str, storage_dir: Path) -> Optional[str]:
    if not raw_path:
        return None
    raw_path = raw_path.strip()
    if raw_path.startswith("storage:"):
        rel_name = raw_path.split("storage:", 1)[1].lstrip("/").strip()
        if not rel_name:
            return None
        full = storage_dir / attachment_key / rel_name
        return str(full)
    if raw_path.startswith("/"):
        return raw_path
    return None


def build_manifest_from_zotero_collection(
    collection: str,
    zotero_db: Path,
    zotero_storage_dir: Path,
) -> Dict:
    if not zotero_db.exists():
        raise ValueError(f"zotero db not found: {zotero_db}")
    if not zotero_storage_dir.exists():
        raise ValueError(f"zotero storage dir not found: {zotero_storage_dir}")

    conn = sqlite3.connect(f"file:{zotero_db}?mode=ro&immutable=1", uri=True)
    conn.row_factory = sqlite3.Row
    cur = conn.cursor()

    collection_row = cur.execute(
        """
        SELECT collectionID, key, collectionName
        FROM collections
        WHERE collectionName = ? OR key = ?
        ORDER BY collectionID
        LIMIT 1
        """,
        (collection, collection),
    ).fetchone()
    if collection_row is None:
        conn.close()
        raise ValueError(f"collection not found in zotero db: {collection}")

    collection_id = int(collection_row["collectionID"])
    collection_name = collection_row["collectionName"] or collection
    collection_key = collection_row["key"] or ""

    parent_rows = cur.execute(
        """
        SELECT DISTINCT ci.itemID, it.typeName
        FROM collectionItems ci
        JOIN items i ON i.itemID = ci.itemID
        JOIN itemTypes it ON it.itemTypeID = i.itemTypeID
        WHERE ci.collectionID = ? AND i.itemID NOT IN (
          SELECT itemID FROM itemAttachments
        )
        ORDER BY ci.itemID
        """,
        (collection_id,),
    ).fetchall()
    parent_ids = [int(r["itemID"]) for r in parent_rows]
    item_type_map: Dict[int, str] = {int(r["itemID"]): str(r["typeName"]) for r in parent_rows}
    if not parent_ids:
        conn.close()
        return {"collection": collection_name, "collection_key": collection_key, "items": []}

    placeholder = ",".join(["?"] * len(parent_ids))
    meta_rows = cur.execute(
        f"""
        SELECT
            id.itemID,
            f.fieldName,
            dv.value
        FROM itemData id
        JOIN fields f ON f.fieldID = id.fieldID
        JOIN itemDataValues dv ON dv.valueID = id.valueID
        WHERE id.itemID IN ({placeholder})
        """,
        parent_ids,
    ).fetchall()

    meta_map: Dict[int, Dict[str, str]] = defaultdict(dict)
    for r in meta_rows:
        meta_map[int(r["itemID"])][str(r["fieldName"])] = str(r["value"])

    creator_rows = cur.execute(
        f"""
        SELECT
            ic.itemID,
            c.firstName,
            c.lastName
        FROM itemCreators ic
        JOIN creators c ON c.creatorID = ic.creatorID
        WHERE ic.itemID IN ({placeholder})
        ORDER BY ic.itemID, ic.orderIndex
        """,
        parent_ids,
    ).fetchall()
    creators_map: Dict[int, List[str]] = defaultdict(list)
    for r in creator_rows:
        first = (r["firstName"] or "").strip()
        last = (r["lastName"] or "").strip()
        full = f"{first} {last}".strip() or last or first
        if full:
            creators_map[int(r["itemID"])].append(full)

    attach_rows = cur.execute(
        f"""
        SELECT
            a.parentItemID,
            a.itemID AS attachmentItemID,
            ai.key AS attachmentKey,
            a.path,
            a.contentType
        FROM itemAttachments a
        JOIN items ai ON ai.itemID = a.itemID
        WHERE a.parentItemID IN ({placeholder})
        ORDER BY a.parentItemID, a.itemID
        """,
        parent_ids,
    ).fetchall()

    attachment_map: Dict[int, Dict] = {}
    for r in attach_rows:
        parent_id = int(r["parentItemID"])
        content_type = (r["contentType"] or "").lower()
        raw_path = r["path"] or ""
        if "pdf" not in content_type and not str(raw_path).lower().endswith(".pdf"):
            continue
        resolved = resolve_pdf_path(raw_path, r["attachmentKey"], zotero_storage_dir)
        if not resolved:
            continue
        if parent_id not in attachment_map:
            attachment_map[parent_id] = {"pdf_path": resolved}

    items: List[Dict] = []
    for item_id in parent_ids:
        attachment = attachment_map.get(item_id)
        if not attachment:
            continue
        meta = meta_map.get(item_id, {})
        title = choose_first_nonempty([meta.get("title")], default="unknown")
        venue = choose_first_nonempty(
            [
                meta.get("publicationTitle"),
                meta.get("proceedingsTitle"),
                meta.get("conferenceName"),
                meta.get("journalAbbreviation"),
                meta.get("repository"),
                meta.get("libraryCatalog"),
            ],
            default="unknown",
        )
        year = parse_year(choose_first_nonempty([meta.get("date"), meta.get("year")], default=""))
        doi = choose_first_nonempty([meta.get("DOI"), meta.get("doi")], default="")
        institution = choose_first_nonempty(
            [meta.get("publisher"), meta.get("university"), meta.get("institution")],
            default="unknown",
        )
        items.append(
            {
                "title": title,
                "authors": creators_map.get(item_id, []),
                "institution": institution,
                "venue": venue,
                "year": year,
                "doi": doi,
                "item_type": item_type_map.get(item_id, "unknown"),
                "library_catalog": choose_first_nonempty([meta.get("libraryCatalog")], default="unknown"),
                "pdf_path": attachment["pdf_path"],
            }
        )

    conn.close()
    return {"collection": collection_name, "collection_key": collection_key, "items": items}


def normalize_title(title: str) -> str:
    return re.sub(r"\s+", " ", re.sub(r"[^a-z0-9 ]", "", title.lower())).strip()


def split_sentences(text: str) -> List[str]:
    text = re.sub(r"\s+", " ", text).strip()
    if not text:
        return []
    parts = re.split(r"(?<=[。！？!?\.])\s+", text)
    return [p.strip() for p in parts if p.strip()]


def extract_keywords_rule(intro_text: str, language: str, top_n: int = 6) -> List[str]:
    text = re.sub(r"\s+", " ", intro_text or "").strip()
    if not text:
        return []
    if language == "zh":
        candidates = re.findall(r"[\u4e00-\u9fff]{2,8}", text)
        stop = {"我们", "本文", "研究", "方法", "结果", "系统", "进行", "通过", "一种", "提出"}
    else:
        candidates = [x.lower() for x in re.findall(r"[A-Za-z][A-Za-z\-]{2,}", text)]
        stop = {
            "this",
            "paper",
            "study",
            "method",
            "results",
            "approach",
            "using",
            "based",
            "propose",
        }
    cnt = Counter(x for x in candidates if x not in stop)
    return [k for k, _ in cnt.most_common(top_n)]


def extract_keywords_from_intro(
    intro_text: str, language: str, llm_client: Optional["LLMClient"], title: str
) -> List[str]:
    intro_text = (intro_text or "")[:2600]
    if llm_client and llm_client.enabled and intro_text.strip():
        prompt_lang = "中文" if language == "zh" else "English"
        system = "You extract concise paper keywords. Return strict JSON only."
        user = (
            f"语言: {prompt_lang}\n"
            "从引言抽取3-8个关键词，输出JSON字段 keywords: string[]。\n"
            f"题目: {title}\n引言:\n{intro_text}\n"
        )
        resp = llm_client.complete_json(system, user)
        if isinstance(resp, dict) and isinstance(resp.get("keywords"), list):
            kws = [str(x).strip() for x in resp.get("keywords", []) if str(x).strip()]
            if kws:
                return kws[:8]
    return extract_keywords_rule(intro_text, language, top_n=6)


def is_noise_sentence(text: str) -> bool:
    t = re.sub(r"\s+", " ", (text or "")).strip().lower()
    if not t:
        return True
    if re.search(r"https?://|www\.", t):
        return True
    if any(x in t for x in ["latest updates", "cookie", "privacy policy", "terms of use", "all rights reserved"]):
        return True
    if len(t) < 15:
        return True
    # OCR/web nav garbage often has many symbols or path-like fragments.
    if re.search(r"[|]{2,}|[/]{2,}|[>]{2,}", t):
        return True
    return False


def truncate_line(text: str, language: str, max_zh_chars: int = 28, max_en_words: int = 18) -> str:
    # Keep full content; only normalize whitespace.
    return re.sub(r"\s+", " ", text).strip()


def extract_text_and_sentences(doc) -> Tuple[List[str], List[SentenceSpan]]:
    page_texts = []
    spans: List[SentenceSpan] = []
    for i, page in enumerate(doc):
        text = page.get_text("text") or ""
        page_texts.append(text)
        for s in split_sentences(text):
            spans.append(SentenceSpan(text=s, page=i + 1))
    return page_texts, spans


def extract_abstract(full_text: str) -> str:
    m = re.search(
        r"(?is)\babstract\b\s*[:\-]?\s*(.*?)\n\s*(?:1\.?\s*introduction|introduction)\b",
        full_text,
    )
    if m:
        return re.sub(r"\s+", " ", m.group(1)).strip()
    return re.sub(r"\s+", " ", full_text[:1800]).strip()


def extract_sections(page_texts: List[str]) -> List[str]:
    sections: List[str] = []
    for text in page_texts:
        for line in text.splitlines():
            s = line.strip()
            if not s:
                continue
            if len(s) > 120:
                continue
            if re.match(r"^(\d+(\.\d+)*)\s+[A-Z].*", s):
                sections.append(s)
            elif re.match(r"^[A-Z][A-Z\s\-]{4,}$", s):
                sections.append(s)
    # unique in order
    seen = set()
    out = []
    for s in sections:
        if s not in seen:
            seen.add(s)
            out.append(s)
    return out[:30]


def extract_captions(page_texts: List[str]) -> List[Dict]:
    out: List[Dict] = []
    for i, text in enumerate(page_texts):
        for line in text.splitlines():
            s = line.strip()
            if re.match(r"^(Figure|Fig\.)\s*\d+", s, re.IGNORECASE):
                out.append({"page": i + 1, "caption": s})
    return out


def extract_best_image(doc, output_dir: Path, stem: str, page_texts: List[str]) -> Optional[str]:
    # Strict mode: only extract Figure 1 / Fig. 1 image.
    image_dir = output_dir / "_images"
    image_dir.mkdir(parents=True, exist_ok=True)
    fig1_pages: List[int] = []
    for i, text in enumerate(page_texts):
        t = (text or "").lower()
        if re.search(r"\b(fig(?:ure)?\.?\s*1)\b", t):
            fig1_pages.append(i)
    for page_idx in fig1_pages:
        page = doc[page_idx]
        images = page.get_images(full=True)
        for img in images:
            xref = img[0]
            try:
                base_image = doc.extract_image(xref)
            except Exception:
                continue
            if not base_image or "image" not in base_image:
                continue
            ext = base_image.get("ext", "png")
            out_path = image_dir / f"{stem}.{ext}"
            out_path.write_bytes(base_image["image"])
            return str(out_path)
    return None


def pick_sentences(spans: List[SentenceSpan], keywords: List[str], limit: int) -> List[SentenceSpan]:
    chosen: List[SentenceSpan] = []
    lower_keywords = [k.lower() for k in keywords]
    for span in spans:
        if is_noise_sentence(span.text):
            continue
        t = span.text.lower()
        if any(k in t for k in lower_keywords):
            chosen.append(span)
        if len(chosen) >= limit:
            break
    return chosen


def fallback_sentences(spans: List[SentenceSpan], limit: int) -> List[SentenceSpan]:
    out = []
    for s in spans:
        if is_noise_sentence(s.text):
            continue
        if 30 <= len(s.text) <= 320:
            out.append(s)
        if len(out) >= limit:
            break
    return out


def detect_open_source(full_text: str) -> Tuple[str, Optional[int]]:
    text = full_text.lower()
    if "github" in text or "open source" in text or "code is available" in text:
        return "yes", 1
    if "code unavailable" in text:
        return "no", 1
    return "unknown", None


def make_items(spans: List[SentenceSpan], language: str, limit: int) -> List[Dict]:
    out = []
    for s in spans[:limit]:
        out.append(
            {
                "text": truncate_line(s.text, language),
                "evidence": [{"page": s.page, "quote": truncate_line(s.text, "en", 120, 35)}],
            }
        )
    return out


def llm_items(raw: Any, language: str, limit: int) -> List[Dict]:
    out: List[Dict] = []
    if not isinstance(raw, list):
        return out
    for x in raw[:limit]:
        if isinstance(x, str):
            text = truncate_line(x.strip(), language)
            page = None
        elif isinstance(x, dict):
            text = truncate_line(
                str(
                    x.get("text")
                    or x.get("content")
                    or x.get("summary")
                    or x.get("point")
                    or ""
                ).strip(),
                language,
            )
            page = x.get("page") or x.get("p")
        else:
            continue
        try:
            page_int = int(page)
        except Exception:
            page_int = None
        if not text:
            continue
        out.append(
            {
                "text": text,
                "evidence": [] if page_int is None else [{"page": page_int, "quote": text}],
            }
        )
    return out


def llm_items_from_raw_text(raw_text: str, language: str, limit: int) -> List[Dict]:
    """Best-effort extractor for truncated/non-JSON-compliant LLM output."""
    out: List[Dict] = []
    if not raw_text:
        return out
    seen = set()
    for m in re.finditer(r'"text"\s*:\s*"((?:\\.|[^"\\])*)"', raw_text):
        s = m.group(1)
        try:
            text = json.loads(f'"{s}"')
        except Exception:
            text = s.replace('\\"', '"').replace("\\n", " ").replace("\\t", " ")
        text = truncate_line(str(text).strip(), language)
        if not text or text in seen or is_noise_sentence(text):
            continue
        seen.add(text)
        out.append({"text": text, "evidence": []})
        if len(out) >= limit:
            break
    if len(out) < limit:
        for m in re.finditer(r"^\s*(?:[-*]|\d+[.)、])\s+(.+?)\s*$", raw_text, re.M):
            text = truncate_line(m.group(1).strip(), language)
            if not text or text in seen or is_noise_sentence(text):
                continue
            seen.add(text)
            out.append({"text": text, "evidence": []})
            if len(out) >= limit:
                break
    if len(out) < limit:
        for line in split_sentences(raw_text):
            text = truncate_line(line.strip(), language)
            if not text or text in seen or is_noise_sentence(text):
                continue
            seen.add(text)
            out.append({"text": text, "evidence": []})
            if len(out) >= limit:
                break
    return out


def detect_section_pages(page_texts: List[str]) -> Dict[str, set]:
    pages: Dict[str, set] = {
        "abstract": set(),
        "intro": set(),
        "method": set(),
        "results": set(),
        "discussion": set(),
        "conclusion": set(),
        "limitations": set(),
    }
    patterns = SECTION_PATTERNS
    for i, text in enumerate(page_texts, start=1):
        for line in text.splitlines():
            s = line.strip().lower()
            if not s or len(s) > 120:
                continue
            is_heading = bool(re.match(r"^(\d+(\.\d+)*)\s+", s)) or s.isupper() or len(s.split()) <= 8
            if not is_heading:
                continue
            for key, pats in patterns.items():
                if any(re.search(p, s) for p in pats):
                    pages[key].add(i)
    if not pages["abstract"] and page_texts:
        # Most papers place abstract on the first page.
        pages["abstract"].add(1)
    return pages


def topic_section_priority(topic: str) -> List[str]:
    # Fixed mapping per user requirement.
    if topic == "task":
        return ["intro"]
    if topic == "method":
        return ["method", "intro"]
    if topic == "contrib":
        return ["intro", "results"]
    return ["conclusion", "limitations", "discussion", "results"]


def extract_section_sentence_spans(page_texts: List[str]) -> Dict[str, List[SentenceSpan]]:
    """
    Parse section blocks by heading boundaries and collect sentence spans per section.
    This avoids mixing Abstract and Introduction on the same page.
    """
    section_map: Dict[str, List[SentenceSpan]] = {
        "abstract": [],
        "intro": [],
        "method": [],
        "results": [],
        "discussion": [],
        "conclusion": [],
        "limitations": [],
    }
    patterns = SECTION_PATTERNS
    current_section: Optional[str] = None

    for page_i, text in enumerate(page_texts, start=1):
        for raw in text.splitlines():
            line = re.sub(r"\s+", " ", raw).strip()
            if not line:
                continue
            low = line.lower()
            if len(low) > 140:
                is_heading = False
            else:
                is_heading = bool(re.match(r"^(\d+(\.\d+)*|[ivx]+\.?)\s+", low)) or len(low.split()) <= 10

            matched_key = None
            if is_heading:
                for key, pats in patterns.items():
                    if any(re.search(p, low) for p in pats):
                        matched_key = key
                        break
            if matched_key:
                current_section = matched_key
                continue

            if not current_section:
                continue
            for s in split_sentences(line):
                if is_noise_sentence(s):
                    continue
                section_map[current_section].append(SentenceSpan(text=s, page=page_i))
    return section_map


def scoped_spans_for_topic(
    spans: List[SentenceSpan],
    page_texts: List[str],
    topic: str,
    section_sentence_map: Optional[Dict[str, List[SentenceSpan]]] = None,
) -> List[SentenceSpan]:
    priority = topic_section_priority(topic)

    # 1) Prefer heading-bounded section spans (paragraph-level).
    if section_sentence_map:
        selected_spans: List[SentenceSpan] = []
        for sec in priority:
            selected_spans.extend(section_sentence_map.get(sec, []))
        if selected_spans:
            return selected_spans

    # 2) Fallback to page-level section detection.
    section_pages = detect_section_pages(page_texts)
    selected_pages = set()
    for sec in priority:
        selected_pages.update(section_pages.get(sec, set()))
    if not selected_pages:
        # For task extraction, do not fall back to full paper.
        if topic == "task":
            return []
        return spans
    return [s for s in spans if s.page in selected_pages]


def build_context_from_spans(spans: List[SentenceSpan], max_lines: int, max_chars: int) -> str:
    lines = []
    total = 0
    for s in spans:
        line = f"[p{s.page}] {re.sub(r'\\s+', ' ', s.text).strip()}"
        if not line.strip():
            continue
        total += len(line) + 1
        if total > max_chars:
            break
        lines.append(line)
        if len(lines) >= max_lines:
            break
    return "\n".join(lines)


def section_context_for_topic(
    spans: List[SentenceSpan],
    page_texts: List[str],
    topic: str,
    keywords: List[str],
    section_sentence_map: Optional[Dict[str, List[SentenceSpan]]] = None,
    max_lines: int = 34,
    max_chars: int = 3200,
) -> str:
    scoped = scoped_spans_for_topic(spans, page_texts, topic, section_sentence_map=section_sentence_map)
    sec_spans = scoped if scoped else []

    picked = pick_sentences(sec_spans, keywords, max_lines) if sec_spans else []
    if len(picked) < max_lines // 2:
        picked += fallback_sentences(sec_spans, max_lines - len(picked))
    if len(picked) < max_lines // 2:
        picked += pick_sentences(spans, keywords, max_lines - len(picked))
    if len(picked) < max_lines // 2:
        picked += fallback_sentences(spans, max_lines - len(picked))

    return build_context_from_spans(picked[:max_lines], max_lines=max_lines, max_chars=max_chars)


def analyze_paper(
    item: Dict,
    output_dir: Path,
    include_images: bool,
    language: str,
    llm_client: Optional[LLMClient] = None,
    progress_cb=None,
) -> Dict:
    pdf_path = item.get("pdf_path", "")
    result = {
        "title": item.get("title") or "unknown",
        "authors": item.get("authors") or [],
        "institution": item.get("institution") or "unknown",
        "venue": item.get("venue") or "unknown",
        "year": item.get("year") or "unknown",
        "doi": item.get("doi") or "",
        "item_type": item.get("item_type") or "unknown",
        "library_catalog": item.get("library_catalog") or "unknown",
        "pdf_path": pdf_path,
        "parse_failed": False,
        "failure_reason": "",
    }
    try:
        if fitz is None:
            raise RuntimeError("PyMuPDF is not installed")
        doc = fitz.open(pdf_path)
        page_texts, spans = extract_text_and_sentences(doc)
        section_sentence_map = extract_section_sentence_spans(page_texts)
        full_text = "\n".join(page_texts)
        abstract = extract_abstract(full_text)
        intro = " ".join(page_texts[:2]).strip()
        sections = extract_sections(page_texts)
        captions = extract_captions(page_texts)

        task_scope = scoped_spans_for_topic(spans, page_texts, "task", section_sentence_map=section_sentence_map)
        method_scope = scoped_spans_for_topic(spans, page_texts, "method", section_sentence_map=section_sentence_map)
        contrib_scope = scoped_spans_for_topic(spans, page_texts, "contrib", section_sentence_map=section_sentence_map)
        limits_scope = scoped_spans_for_topic(spans, page_texts, "limits", section_sentence_map=section_sentence_map)

        task_spans = pick_sentences(task_scope, ["problem", "task", "we address", "we study", "aim", "this paper", "we propose"], 1)
        if not task_spans:
            task_spans = fallback_sentences(task_scope, 1)

        method_spans = pick_sentences(
            method_scope,
            ["we propose", "framework", "method", "model", "pipeline", "approach"],
            3,
        )
        if not method_spans:
            method_spans = fallback_sentences(method_scope or spans, 3)

        contrib_spans = pick_sentences(
            contrib_scope,
            ["contribution", "we show", "we demonstrate", "outperform", "results"],
            4,
        )
        if not contrib_spans:
            contrib_spans = fallback_sentences(contrib_scope or spans, 4)

        limit_spans = pick_sentences(
            limits_scope,
            ["limitation", "future work", "challenge", "fails", "bottleneck"],
            3,
        )
        if not limit_spans:
            scoped_tail = limits_scope[-12:] if len(limits_scope) > 12 else limits_scope
            limit_spans = fallback_sentences(scoped_tail or (spans[-12:] if len(spans) > 12 else spans), 2)

        open_source_status, open_source_page = detect_open_source(full_text)
        keywords = extract_keywords_from_intro(intro, language, llm_client, result["title"])

        llm_method = []
        llm_contrib = []
        llm_limits = []
        llm_task = []
        llm_open_source_status = None
        llm_task_used = False
        llm_method_used = False
        llm_contrib_used = False
        llm_limits_used = False
        llm_error = ""
        llm_error_debug = ""
        llm_raw_preview = ""
        if llm_client and llm_client.enabled:
            if progress_cb:
                progress_cb("llm:prepare")
            prompt_lang = "中文" if language == "zh" else "English"
            llm_system = "You are an expert research summarizer. Return strict JSON only."
            errors = []
            previews = []

            task_ctx = section_context_for_topic(
                spans,
                page_texts,
                topic="task",
                keywords=["problem", "task", "we address", "we study", "aim", "this paper", "we propose"],
                section_sentence_map=section_sentence_map,
                max_lines=14,
                max_chars=min(llm_client.max_input_chars, 1200),
            )
            task_user = (
                f"语言: {prompt_lang}\n"
                "输出JSON字段：task_definition:[{text,page}]，严格1条，中文优先，一句话概括研究任务。\n"
                f"题目: {result['title']}\n证据:\n{task_ctx}\n"
            )
            if progress_cb:
                progress_cb("llm:task:start")
            resp_t = llm_client.complete_json(llm_system, task_user)
            if progress_cb:
                progress_cb(f"llm:task:done used={bool(resp_t)}")
            if resp_t:
                task_raw = (
                    resp_t.get("task_definition")
                    or resp_t.get("task")
                    or resp_t.get("problem_definition")
                    or resp_t.get("研究任务")
                    or resp_t.get("任务定义")
                )
                llm_task = llm_items(task_raw, language, 1)
                llm_task_used = len(llm_task) > 0
            else:
                llm_task = llm_items_from_raw_text(llm_client.last_raw_content or "", language, 1)
                if not llm_task:
                    text_fallback = llm_client.complete_text(
                        llm_system,
                        (
                            f"语言: {prompt_lang}\n"
                            "请仅输出1条任务定义，一行，不要解释，不要标题。\n"
                            f"题目: {result['title']}\n证据:\n{task_ctx}\n"
                        ),
                    )
                    llm_task = llm_items_from_raw_text(text_fallback or "", language, 1)
                llm_task_used = len(llm_task) > 0
                errors.append(f"task:{llm_client.last_error}")
                previews.append((llm_client.last_raw_content or "")[:140])

            method_ctx = section_context_for_topic(
                spans,
                page_texts,
                topic="method",
                keywords=["we propose", "framework", "method", "pipeline", "approach", "architecture"],
                section_sentence_map=section_sentence_map,
                max_lines=22,
                max_chars=min(llm_client.max_input_chars, 1800),
            )
            method_user = (
                f"语言: {prompt_lang}\n"
                "输出JSON字段：core_method:[{text,page}]，2-3条，中文优先，短句。\n"
                f"题目: {result['title']}\n证据:\n{method_ctx}\n"
            )
            if progress_cb:
                progress_cb("llm:method:start")
            resp_m = llm_client.complete_json(llm_system, method_user)
            if progress_cb:
                progress_cb(f"llm:method:done used={bool(resp_m)}")
            if resp_m:
                llm_method = llm_items(resp_m.get("core_method"), language, 3)
                llm_method_used = len(llm_method) > 0
            else:
                llm_method = llm_items_from_raw_text(llm_client.last_raw_content or "", language, 3)
                if not llm_method:
                    text_fallback = llm_client.complete_text(
                        llm_system,
                        (
                            f"语言: {prompt_lang}\n"
                            "请仅输出2-3条核心方法要点，每行一条，不要解释，不要标题。\n"
                            f"题目: {result['title']}\n证据:\n{method_ctx}\n"
                        ),
                    )
                    llm_method = llm_items_from_raw_text(text_fallback or "", language, 3)
                llm_method_used = len(llm_method) > 0
                errors.append(f"method:{llm_client.last_error}")
                previews.append((llm_client.last_raw_content or "")[:140])

            contrib_ctx = section_context_for_topic(
                spans,
                page_texts,
                topic="contrib",
                keywords=["contribution", "results", "outperform", "we show", "we demonstrate", "experiment"],
                section_sentence_map=section_sentence_map,
                max_lines=26,
                max_chars=min(llm_client.max_input_chars, 2200),
            )
            contrib_user = (
                f"语言: {prompt_lang}\n"
                "输出JSON字段：main_contributions:[{text,page}]，2-4条，中文优先，强调创新点与实验结论。\n"
                f"题目: {result['title']}\n证据:\n{contrib_ctx}\n"
            )
            if progress_cb:
                progress_cb("llm:contrib:start")
            resp_c = llm_client.complete_json(llm_system, contrib_user)
            if progress_cb:
                progress_cb(f"llm:contrib:done used={bool(resp_c)}")
            if resp_c:
                contrib_raw = (
                    resp_c.get("main_contributions")
                    or resp_c.get("contributions")
                    or resp_c.get("major_contributions")
                    or resp_c.get("highlights")
                    or resp_c.get("关键贡献")
                    or resp_c.get("主要贡献")
                    or resp_c.get("贡献")
                )
                llm_contrib = llm_items(contrib_raw, language, 4)
                llm_contrib_used = len(llm_contrib) > 0
                o = str(resp_c.get("open_source_status", "")).strip().lower()
                if o in {"yes", "no", "unknown"}:
                    llm_open_source_status = o
            else:
                llm_contrib = llm_items_from_raw_text(llm_client.last_raw_content or "", language, 4)
                if not llm_contrib:
                    text_fallback = llm_client.complete_text(
                        llm_system,
                        (
                            f"语言: {prompt_lang}\n"
                            "请仅输出2-4条要点，每行一条，不要解释，不要编号标题。\n"
                            "主题：论文主要贡献（创新点+关键实验结论）。\n"
                            f"题目: {result['title']}\n证据:\n{contrib_ctx}\n"
                        ),
                    )
                    llm_contrib = llm_items_from_raw_text(text_fallback or "", language, 4)
                llm_contrib_used = len(llm_contrib) > 0
                errors.append(f"contrib:{llm_client.last_error}")
                previews.append((llm_client.last_raw_content or "")[:140])

            limits_ctx = section_context_for_topic(
                spans,
                page_texts,
                topic="limits",
                keywords=["limitation", "future work", "challenge", "bottleneck", "fails"],
                section_sentence_map=section_sentence_map,
                max_lines=20,
                max_chars=min(llm_client.max_input_chars, 1600),
            )
            limits_user = (
                f"语言: {prompt_lang}\n"
                "输出JSON字段：limitations:[{text,page}]，1-3条，中文优先，聚焦不足与边界。\n"
                f"题目: {result['title']}\n证据:\n{limits_ctx}\n"
            )
            if progress_cb:
                progress_cb("llm:limits:start")
            resp_l = llm_client.complete_json(llm_system, limits_user)
            if progress_cb:
                progress_cb(f"llm:limits:done used={bool(resp_l)}")
            if resp_l:
                llm_limits = llm_items(resp_l.get("limitations"), language, 3)
                llm_limits_used = len(llm_limits) > 0
            else:
                llm_limits = llm_items_from_raw_text(llm_client.last_raw_content or "", language, 3)
                if not llm_limits:
                    text_fallback = llm_client.complete_text(
                        llm_system,
                        (
                            f"语言: {prompt_lang}\n"
                            "请仅输出1-3条局限性，每行一条，不要解释，不要标题。\n"
                            f"题目: {result['title']}\n证据:\n{limits_ctx}\n"
                        ),
                    )
                    llm_limits = llm_items_from_raw_text(text_fallback or "", language, 3)
                llm_limits_used = len(llm_limits) > 0
                errors.append(f"limits:{llm_client.last_error}")
                previews.append((llm_client.last_raw_content or "")[:140])

            # Report only unresolved errors; suppress transient parse/empty errors
            # if fallback extraction eventually succeeded.
            unresolved: List[str] = []
            for e in errors:
                if e.startswith("task:") and llm_task_used:
                    continue
                if e.startswith("method:") and llm_method_used:
                    continue
                if e.startswith("contrib:") and llm_contrib_used:
                    continue
                if e.startswith("limits:") and llm_limits_used:
                    continue
                unresolved.append(e)
            llm_error_debug = "; ".join([e for e in errors if e])
            llm_error = "; ".join([e for e in unresolved if e])
            llm_raw_preview = " | ".join([p for p in previews if p])[:320]

        image_path = None
        if include_images:
            if progress_cb:
                progress_cb("image:extract:start")
            image_path = extract_best_image(
                doc,
                output_dir,
                normalize_title(result["title"])[:80] or "paper",
                page_texts,
            )
            if progress_cb:
                progress_cb(f"image:extract:done has_image={bool(image_path)}")

        result.update(
            {
                "abstract": abstract,
                "intro": intro,
                "sections": sections,
                "figure_captions": captions,
                "keywords": keywords,
                "task_definition": llm_task or make_items(task_spans, language, 1),
                "core_method": llm_method or make_items(method_spans, language, 3),
                "main_contributions": llm_contrib or make_items(contrib_spans, language, 4),
                "limitations": llm_limits or make_items(limit_spans, language, 3),
                "open_source_status": {
                    "value": llm_open_source_status if llm_open_source_status in {"yes", "no", "unknown"} else open_source_status,
                    "evidence": [] if open_source_page is None else [{"page": open_source_page, "quote": "open-source signal"}],
                },
                "llm_task_used": llm_task_used,
                "llm_method_used": llm_method_used,
                "llm_contrib_used": llm_contrib_used,
                "llm_limits_used": llm_limits_used,
                "llm_error": llm_error,
                "llm_error_debug": llm_error_debug,
                "llm_raw_preview": llm_raw_preview,
                "image_path": image_path,
            }
        )

        doc.close()
        return result
    except Exception as e:
        result["parse_failed"] = True
        result["failure_reason"] = str(e)
        result.update(
            {
                "abstract": "",
                "intro": "",
                "sections": [],
                "figure_captions": [],
                "keywords": [],
                "task_definition": [],
                "core_method": [],
                "main_contributions": [],
                "limitations": [],
                "open_source_status": {"value": "unknown", "evidence": []},
                "llm_task_used": False,
                "llm_method_used": False,
                "llm_contrib_used": False,
                "llm_limits_used": False,
                "llm_error": "",
                "llm_error_debug": "",
                "llm_raw_preview": "",
                "image_path": None,
            }
        )
        return result


def dedup_items(items: List[Dict]) -> List[Dict]:
    def score(it: Dict) -> int:
        s = 0
        if (it.get("title") or "").strip() and it.get("title") != "unknown":
            s += 1
        if (it.get("year") or "").strip() and it.get("year") != "unknown":
            s += 1
        if isinstance(it.get("authors"), list) and len(it.get("authors")) > 0:
            s += 2
        if (it.get("venue") or "").strip() and it.get("venue") != "unknown":
            s += 3
        if (it.get("institution") or "").strip() and it.get("institution") != "unknown":
            s += 2
        if (it.get("library_catalog") or "").strip() and it.get("library_catalog") != "unknown":
            s += 1
        if (it.get("doi") or "").strip():
            s += 1
        if (it.get("pdf_path") or "").strip():
            s += 1
        return s

    best_by_key: Dict[Tuple[str, str], Dict] = {}
    for item in items:
        doi = (item.get("doi") or "").strip().lower()
        title = normalize_title(item.get("title") or "")
        key = ("doi", doi) if doi else ("title", title)
        prev = best_by_key.get(key)
        if prev is None or score(item) > score(prev):
            best_by_key[key] = item
    return list(best_by_key.values())


def as_doc_text(p: Dict) -> str:
    bits = []
    for k in ["task_definition", "core_method", "main_contributions", "limitations"]:
        bits.extend([x["text"] for x in p.get(k, [])])
    bits.extend(p.get("sections", [])[:8])
    return " ".join(bits).strip()


def cluster_papers(papers: List[Dict], forced_k: int = 0) -> Dict[int, List[int]]:
    n = len(papers)
    if n == 0:
        return {}
    if n == 1 or KMeans is None or TfidfVectorizer is None:
        return {0: [0]}

    docs = [as_doc_text(p) for p in papers]
    vectorizer = TfidfVectorizer(analyzer="char_wb", ngram_range=(2, 4), min_df=1)
    X = vectorizer.fit_transform(docs)

    max_k = min(5, n)
    if forced_k > 0:
        k = max(1, min(forced_k, n))
        model = KMeans(n_clusters=k, n_init=10, random_state=SEED)
        labels = model.fit_predict(X)
        clusters: Dict[int, List[int]] = defaultdict(list)
        for i, label in enumerate(labels):
            clusters[int(label)].append(i)
        return dict(sorted(clusters.items(), key=lambda kv: kv[0]))
    best_k = 2 if n >= 2 else 1
    best_score = -1.0

    for k in range(2, max_k + 1):
        model = KMeans(n_clusters=k, n_init=10, random_state=SEED)
        labels = model.fit_predict(X)
        if len(set(labels)) < 2:
            continue
        try:
            score = silhouette_score(X, labels)
        except Exception:
            score = -1.0
        if score > best_score:
            best_score = score
            best_k = k

    model = KMeans(n_clusters=best_k, n_init=10, random_state=SEED)
    labels = model.fit_predict(X)

    clusters: Dict[int, List[int]] = defaultdict(list)
    for i, label in enumerate(labels):
        clusters[int(label)].append(i)
    return dict(sorted(clusters.items(), key=lambda kv: kv[0]))


def compose_cluster_summary(cluster_papers_: List[Dict], language: str) -> Dict:
    tasks = [x["text"] for p in cluster_papers_ for x in p.get("task_definition", [])]
    methods = [x["text"] for p in cluster_papers_ for x in p.get("core_method", [])]
    limits = [x["text"] for p in cluster_papers_ for x in p.get("limitations", [])]

    def pick(lines: List[str], fallback: str, limit: int = 90) -> str:
        if not lines:
            return fallback
        c = Counter(lines)
        text = c.most_common(1)[0][0]
        if language == "zh" and len(text) > limit:
            return text[: limit - 1] + "…"
        return text

    if language == "zh":
        return {
            "core_problem": pick(tasks, "围绕相近任务设定与约束进行建模。"),
            "main_tech_path": pick(methods, "通过模型/框架设计完成任务求解。"),
            "typical_limits": pick(limits, "在泛化性、鲁棒性或数据依赖上存在不足。"),
        }
    return {
        "core_problem": pick(tasks, "Solve related task settings under common constraints."),
        "main_tech_path": pick(methods, "Use model or framework design as the main approach."),
        "typical_limits": pick(limits, "Limitations remain in generalization, robustness, or data dependence."),
    }


def build_cluster_rag_query(cluster_papers_: List[Dict], language: str) -> str:
    titles = [str(p.get("title", "")).strip() for p in cluster_papers_ if str(p.get("title", "")).strip()]
    kws: List[str] = []
    for p in cluster_papers_:
        for k in p.get("keywords", []) or []:
            kk = str(k).strip()
            if kk:
                kws.append(kk)
    uniq_kws = []
    seen = set()
    for k in kws:
        low = k.lower()
        if low in seen:
            continue
        seen.add(low)
        uniq_kws.append(k)
    title_part = "; ".join(titles[:3])
    kw_part = ", ".join(uniq_kws[:12])
    if language == "zh":
        return f"主题综述检索：{title_part}。关键词：{kw_part}"
    return f"Literature synthesis query: {title_part}. Keywords: {kw_part}"


def format_rag_hits_for_prompt(hits: List[Dict[str, Any]], language: str, max_hits: int = 6) -> str:
    if not hits:
        return ""
    lines = []
    for h in hits[:max_hits]:
        title = str(h.get("title", "")).strip() or "unknown"
        snippet = truncate_line(str(h.get("snippet", "")).strip(), language, 100, 48)
        score = h.get("score")
        score_text = f"{score:.4f}" if isinstance(score, float) else "na"
        lines.append(f"- {title} | score={score_text} | evidence={snippet}")
    return "\n".join(lines)


def make_overview(
    clusters: Dict[int, List[int]],
    papers: List[Dict],
    language: str,
    llm_client: Optional[LLMClient] = None,
    rag_client: Optional[ZoteroMcpRagClient] = None,
    rag_top_k: int = 8,
) -> List[Dict]:
    out = []
    for cid, idxs in clusters.items():
        cluster_set = [papers[i] for i in idxs]
        s = compose_cluster_summary(cluster_set, language)
        rag_query = ""
        rag_hits: List[Dict[str, Any]] = []
        if rag_client and rag_client.enabled:
            rag_query = build_cluster_rag_query(cluster_set, language)
            rag_hits = rag_client.search(rag_query, rag_top_k)
        if llm_client and llm_client.enabled:
            prompt_lang = "中文" if language == "zh" else "English"
            cluster_briefs = []
            for p in cluster_set:
                cluster_briefs.append(
                    {
                        "title": p.get("title", "unknown"),
                        "task_definition": [x.get("text") for x in p.get("task_definition", [])],
                        "core_method": [x.get("text") for x in p.get("core_method", [])],
                        "limitations": [x.get("text") for x in p.get("limitations", [])],
                    }
                )
            llm_system = "You synthesize research directions across multiple papers. Return strict JSON only."
            llm_user = (
                f"语言: {prompt_lang}\n"
                "基于以下同一聚类内论文，做跨论文抽象，输出JSON字段：\n"
                "core_problem, main_tech_path, typical_limits。\n"
                "禁止逐篇拼接，不要长段落，每项一句。\n"
                f"cluster_papers={json.dumps(cluster_briefs, ensure_ascii=False)}"
            )
            rag_text = format_rag_hits_for_prompt(rag_hits, language, max_hits=min(6, rag_top_k))
            if rag_text:
                llm_user += f"\n补充RAG证据（优先参考但需保持跨论文抽象）：\n{rag_text}"
            llm_resp = llm_client.complete_json(llm_system, llm_user)
            if isinstance(llm_resp, dict):
                s = {
                    "core_problem": truncate_line(str(llm_resp.get("core_problem", s["core_problem"])), language, 90, 45),
                    "main_tech_path": truncate_line(str(llm_resp.get("main_tech_path", s["main_tech_path"])), language, 90, 45),
                    "typical_limits": truncate_line(str(llm_resp.get("typical_limits", s["typical_limits"])), language, 90, 45),
                }
        out.append(
            {
                "cluster_id": cid,
                "paper_count": len(cluster_set),
                "core_problem": s["core_problem"],
                "main_tech_path": s["main_tech_path"],
                "typical_limits": s["typical_limits"],
                "rag_query": rag_query,
                "rag_hits": rag_hits[:rag_top_k],
            }
        )
    return out


def make_final_synthesis(
    clusters_summary: List[Dict],
    language: str,
    llm_client: Optional[LLMClient] = None,
    rag_enabled: bool = False,
) -> Dict:
    if language == "zh":
        core = "现有工作围绕任务可行性与泛化能力展开，主线集中在表示设计、学习策略与数据机制。"
        gaps = [
            "跨场景泛化缺乏统一评估协议。",
            "多模态信号融合的可解释性仍不足。",
            "真实部署条件下的鲁棒性验证不充分。",
        ]
    else:
        core = "Existing work focuses on feasibility and generalization, mainly via representation, learning strategy, and data design."
        gaps = [
            "No unified protocol for cross-scenario generalization.",
            "Interpretability in multi-modal fusion remains limited.",
            "Real-world robustness validation is still insufficient.",
        ]
    directions = []
    for c in clusters_summary:
        directions.append(
            {
                "cluster_id": c["cluster_id"],
                "topic": c["core_problem"],
                "advantages": c["main_tech_path"],
                "limitations": c["typical_limits"],
                "potential_gap": c["typical_limits"],
            }
        )
    result = {
        "core_problem_summary": core,
        "directions": directions,
        "possible_research_gaps": gaps[:6],
    }
    if llm_client and llm_client.enabled:
        prompt_lang = "中文" if language == "zh" else "English"
        llm_system = "You generate concise final synthesis for literature review slides. Return strict JSON only."
        llm_user = (
            f"语言: {prompt_lang}\n"
            "基于聚类总结，输出JSON：\n"
            "core_problem_summary: string\n"
            "possible_research_gaps: string[] (3-6)\n"
            "要求：跨论文归纳，避免空话，适合组会PPT。\n"
            f"clusters_summary={json.dumps(clusters_summary, ensure_ascii=False)}"
        )
        if rag_enabled:
            rag_briefs = []
            for c in clusters_summary:
                hits = c.get("rag_hits", [])
                if not isinstance(hits, list):
                    continue
                rag_briefs.append(
                    {
                        "cluster_id": c.get("cluster_id"),
                        "evidence": [
                            {
                                "title": h.get("title", ""),
                                "snippet": truncate_line(str(h.get("snippet", "")), language, 90, 45),
                            }
                            for h in hits[:4]
                            if isinstance(h, dict)
                        ],
                    }
                )
            if rag_briefs:
                llm_user += f"\n补充RAG证据（用于支撑gap归纳）：{json.dumps(rag_briefs, ensure_ascii=False)}"
        llm_resp = llm_client.complete_json(llm_system, llm_user)
        if isinstance(llm_resp, dict):
            maybe_gaps = llm_resp.get("possible_research_gaps")
            if isinstance(maybe_gaps, list):
                cleaned = [truncate_line(str(x), language, 45, 20) for x in maybe_gaps if str(x).strip()]
                if cleaned:
                    result["possible_research_gaps"] = cleaned[:6]
            maybe_core = str(llm_resp.get("core_problem_summary", "")).strip()
            if maybe_core:
                result["core_problem_summary"] = truncate_line(maybe_core, language, 120, 65)
    return result


def overlap_too_high(source: str, target: str) -> bool:
    source = re.sub(r"\s+", " ", source).strip().lower()
    target = re.sub(r"\s+", " ", target).strip().lower()
    if not source or not target:
        return False
    if len(target) < 40:
        return False
    return target in source


def run_quality_gates(papers: List[Dict], overview: List[Dict], final_syn: Dict) -> Dict:
    res = {"paper_checks": [], "overview_checks": [], "final_checks": []}
    for p in papers:
        checks = {"title": p.get("title", "unknown"), "passed": True, "issues": []}
        for section in ["task_definition", "core_method", "main_contributions", "limitations"]:
            vals = p.get(section, [])
            if not vals:
                checks["passed"] = False
                checks["issues"].append(f"missing_{section}")
                continue
            for v in vals:
                if not v.get("evidence"):
                    checks["passed"] = False
                    checks["issues"].append(f"missing_evidence_{section}")
                if overlap_too_high(p.get("abstract", ""), v.get("text", "")):
                    checks["passed"] = False
                    checks["issues"].append(f"abstract_overlap_{section}")
        res["paper_checks"].append(checks)

    if not overview:
        res["overview_checks"].append({"passed": False, "issue": "empty_overview"})
    else:
        res["overview_checks"].append({"passed": True})

    if not final_syn.get("possible_research_gaps"):
        res["final_checks"].append({"passed": False, "issue": "empty_research_gaps"})
    else:
        res["final_checks"].append({"passed": True})
    return res


def zh_num_label(i: int) -> str:
    m = {
        1: "一",
        2: "二",
        3: "三",
        4: "四",
        5: "五",
        6: "六",
        7: "七",
        8: "八",
        9: "九",
        10: "十",
    }
    return m.get(i, str(i))


def direction_title(i: int, topic: str, language: str) -> str:
    t = (topic or "").strip()
    if language == "zh":
        prefix = f"方向{zh_num_label(i)}"
        return f"{prefix}：{t}" if t else prefix
    prefix = f"Direction {i}"
    return f"{prefix}: {t}" if t else prefix


def render_markdown(output: Path, collection: str, overview: List[Dict], papers: List[Dict], final_syn: Dict) -> None:
    lines = []
    lines.append(f"# review_{collection}")
    lines.append("")
    lines.append("## 文献整体结构梳理")
    lines.append("")
    for i, c in enumerate(overview, start=1):
        lines.append(f"### {direction_title(i, c.get('core_problem', ''), 'zh')}")
        lines.append(f"- 核心问题: {c['core_problem']}")
        lines.append(f"- 主流技术路径: {c['main_tech_path']}")
        lines.append(f"- 典型局限: {c['typical_limits']}")
        lines.append("")

    lines.append("## 单篇论文分析")
    lines.append("")
    for p in papers:
        lines.append(f"### {p.get('title', 'unknown')}")
        authors = ", ".join(p.get("authors", [])) if p.get("authors") else "未在Zotero填写"
        lines.append(f"- 作者: {authors}")
        lines.append(f"- 机构: {display_meta(p.get('institution'))}")
        venue = display_meta(p.get("venue"))
        year = display_meta(p.get("year"), fallback="年份缺失")
        if venue == "未在Zotero填写":
            venue = item_type_label(display_meta(p.get("item_type"), fallback="文献类型未知"))
        lines.append(f"- 来源/年份: {venue} {year}")
        td = p.get("task_definition", [])
        cm = p.get("core_method", [])
        mc = p.get("main_contributions", [])
        lm = p.get("limitations", [])
        lines.append("- 任务定义:")
        for x in td[:1]:
            lines.append(f"  - {x['text']}")
        lines.append("- 核心方法:")
        for x in cm[:3]:
            lines.append(f"  - {x['text']}")
        lines.append("- 主要贡献:")
        for x in mc[:4]:
            lines.append(f"  - {x['text']}")
        lines.append("- 局限:")
        for x in lm[:3]:
            lines.append(f"  - {x['text']}")
        lines.append(f"- 是否开源: {p.get('open_source_status', {}).get('value', 'unknown')}")
        if p.get("image_path"):
            lines.append(f"- 图像: {p['image_path']}")
        else:
            lines.append("- 图像: 图像未提取到")
        lines.append("")

    lines.append("## 结构归纳与研究机会")
    lines.append("")
    lines.append(final_syn.get("core_problem_summary", ""))
    lines.append("")
    dirs = final_syn.get("directions", [])
    if isinstance(dirs, list) and dirs:
        lines.append("### 主要技术方向")
        for d in dirs[:5]:
            cid = int(d.get("cluster_id", 0)) + 1
            topic = str(d.get("topic", "")).strip()
            lines.append(f"- {direction_title(cid, topic, 'zh')}")
        lines.append("")
    lines.append("### Possible Research Gaps")
    for g in final_syn.get("possible_research_gaps", []):
        lines.append(f"- {g}")

    output.write_text("\n".join(lines) + "\n", encoding="utf-8")


def add_textbox(slide, x, y, w, h, text, font_size=18, bold=False, color=None, line_spacing=1.0):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.02)
    tf.margin_right = Inches(0.02)
    tf.margin_top = Inches(0.01)
    tf.margin_bottom = Inches(0.01)
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    p.line_spacing = line_spacing
    if color is not None:
        run.font.color.rgb = color
    return tb


def display_meta(v: Any, fallback: str = "未在Zotero填写") -> str:
    text = str(v).strip() if v is not None else ""
    if not text or text.lower() == "unknown":
        return fallback
    return text


def item_type_label(item_type: str) -> str:
    m = {
        "journalArticle": "期刊论文",
        "conferencePaper": "会议论文",
        "book": "书籍",
        "bookSection": "书籍章节",
        "thesis": "学位论文",
        "report": "报告",
        "preprint": "预印本",
        "webpage": "网页资料",
    }
    return m.get(item_type, item_type)


def shorten_for_box(text: str, language: str, hard_zh: int = 34, hard_en_words: int = 16) -> str:
    # No truncation policy: keep full sentence and let text boxes wrap.
    return truncate_line(text, language, max_zh_chars=hard_zh, max_en_words=hard_en_words)


def section_bullets_text(entries: List[Dict], max_n: int, language: str) -> str:
    lines: List[str] = []
    for entry in entries[:max_n]:
        line = shorten_for_box(entry.get("text", ""), language, hard_zh=60, hard_en_words=24)
        if line:
            lines.append(f"- {line}")
    return "\n".join(lines)


def pick_figure_caption(captions: List[str], language: str) -> str:
    if not captions:
        return ""
    for c in captions:
        text = re.sub(r"\s+", " ", str(c)).strip()
        if re.search(r"\b(fig(?:ure)?\.?\s*1|图\s*1)\b", text, re.I):
            return truncate_line(text, language, 68, 30)
    first = re.sub(r"\s+", " ", str(captions[0])).strip()
    return truncate_line(first, language, 68, 30)


def render_ppt(output: Path, overview: List[Dict], papers: List[Dict], final_syn: Dict, language: str) -> Optional[str]:
    if Presentation is None:
        return "python-pptx not installed; skipped pptx output"

    prs = Presentation()
    # 16:9
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blue = RGBColor(0, 162, 232)
    black = RGBColor(0, 0, 0)

    # Part 1
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_textbox(slide, 0.5, 0.2, 12.5, 0.8, "文献整体结构梳理" if language == "zh" else "Literature Structure Overview", 30, True, black)
    y = 1.0
    for i, c in enumerate(overview, start=1):
        add_textbox(slide, 0.6, y, 12.0, 0.5, direction_title(i, c.get("core_problem", ""), language), 20, True, blue)
        y += 0.45
        add_textbox(slide, 0.9, y, 11.8, 0.45, f"核心问题: {c['core_problem']}" if language == "zh" else f"Core Problem: {c['core_problem']}", 14, False, black)
        y += 0.4
        add_textbox(slide, 0.9, y, 11.8, 0.45, f"技术路径: {c['main_tech_path']}" if language == "zh" else f"Tech Path: {c['main_tech_path']}", 14, False, black)
        y += 0.4
        add_textbox(slide, 0.9, y, 11.8, 0.45, f"局限: {c['typical_limits']}" if language == "zh" else f"Limitations: {c['typical_limits']}", 14, False, black)
        y += 0.6
        if y > 6.5:
            break

    # Part 2
    for p in papers:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        title_text = shorten_for_box(p.get("title", "unknown"), language, hard_zh=56, hard_en_words=20)
        add_textbox(slide, 0.5, 0.2, 12.3, 0.85, title_text, 20, True, black)
        authors = ", ".join(p.get("authors", [])) if p.get("authors") else "未在Zotero填写"
        add_textbox(slide, 0.5, 0.78, 12.3, 0.45, authors, 13, False, RGBColor(55, 55, 55))
        venue = display_meta(p.get("venue"))
        year = display_meta(p.get("year"), fallback="年份缺失")
        institution = display_meta(p.get("institution"))
        if venue == "未在Zotero填写":
            venue = item_type_label(display_meta(p.get("item_type"), fallback="文献类型未知"))
        open_src = display_meta(p.get("open_source_status", {}).get("value", "unknown"), fallback="未知")
        inst_text = institution if institution != "未在Zotero填写" else "机构信息缺失"
        src_text = open_src if open_src in {"yes", "no", "unknown", "未知"} else open_src
        meta = f"{venue}（{year}）| {inst_text} | 开源: {src_text}"
        add_textbox(slide, 0.5, 1.1, 7.0, 0.45, meta, 15, False, RGBColor(70, 70, 70))
        doi = (p.get("doi") or "").strip()
        if doi:
            add_textbox(slide, 0.5, 1.45, 6.8, 0.30, f"DOI: {doi}", 11, False, RGBColor(90, 90, 90))
        add_textbox(slide, 7.0, 1.0, 5.8, 0.5, "任务定义" if language == "zh" else "Task", 20, True, blue)

        y_right = 1.45
        for section_title, section_key, max_n in [
            ("任务定义", "task_definition", 1),
            ("核心方法", "core_method", 3),
            ("主要贡献", "main_contributions", 4),
        ]:
            add_textbox(slide, 7.0, y_right, 5.8, 0.35, section_title if language == "zh" else section_title, 15, True, black)
            y_right += 0.24
            block = section_bullets_text(p.get(section_key, []), max_n, language)
            block_h = 0.45 + 0.35 * max(1, len([x for x in p.get(section_key, [])[:max_n] if x.get("text")]))
            add_textbox(slide, 7.2, y_right, 5.5, block_h, block, 13, False, black, line_spacing=1.35)
            y_right += block_h + 0.04

        add_textbox(slide, 0.6, 1.8, 6.0, 0.35, "图像" if language == "zh" else "Image", 15, True, black)
        if p.get("image_path") and Path(p["image_path"]).exists():
            slide.shapes.add_picture(str(p["image_path"]), Inches(0.6), Inches(2.1), width=Inches(5.8), height=Inches(3.4))
        else:
            add_textbox(slide, 0.8, 3.25, 5.2, 0.5, "图像未提取到" if language == "zh" else "Image not extracted", 14, False, RGBColor(100, 100, 100))
        fig_caption = pick_figure_caption(p.get("figure_captions", []), language)
        if fig_caption:
            add_textbox(slide, 0.6, 5.52, 6.0, 0.28, fig_caption, 10, False, RGBColor(120, 120, 120))

        add_textbox(slide, 0.6, 5.88, 6.0, 0.35, "局限" if language == "zh" else "Limitations", 15, True, black)
        y_l = 6.18
        lim_block = section_bullets_text(p.get("limitations", []), 3, language)
        add_textbox(slide, 0.8, y_l, 5.8, 1.0, lim_block, 12, False, black, line_spacing=1.35)

    # Part 3
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_textbox(slide, 0.5, 0.2, 12.2, 0.8, "结构归纳与研究机会" if language == "zh" else "Synthesis and Research Opportunities", 30, True, black)
    add_textbox(slide, 0.7, 1.1, 12.0, 0.9, final_syn.get("core_problem_summary", ""), 16, False, black)

    y = 2.1
    for d in final_syn.get("directions", [])[:5]:
        topic = str(d.get("topic", "")).strip()
        add_textbox(
            slide,
            0.7,
            y,
            12.0,
            0.35,
            direction_title(int(d["cluster_id"]) + 1, topic, language),
            15,
            True,
            blue,
        )
        y += 0.3
        add_textbox(slide, 1.0, y, 11.7, 0.3, ("优势: " if language == "zh" else "Advantage: ") + d["advantages"], 12, False, black)
        y += 0.25
        add_textbox(slide, 1.0, y, 11.7, 0.3, ("局限: " if language == "zh" else "Limitation: ") + d["limitations"], 12, False, black)
        y += 0.38
        if y > 5.5:
            break

    add_textbox(slide, 0.7, 5.8, 12.0, 0.4, "Possible Research Gaps", 16, True, black)
    y = 6.15
    for g in final_syn.get("possible_research_gaps", [])[:6]:
        add_textbox(slide, 0.9, y, 11.8, 0.3, f"- {g}", 12, False, black)
        y += 0.25

    prs.save(str(output))
    return None


def get_stage_paths(args: argparse.Namespace, output_dir: Path, stem: str) -> Dict[str, Path]:
    analyze_json = Path(args.analyze_json) if args.analyze_json else output_dir / f"{stem}.analyze.json"
    global_json = Path(args.global_json) if args.global_json else output_dir / f"{stem}.global.json"
    report_json = output_dir / f"{stem}.json"
    out_md = output_dir / f"{stem}.md"
    out_pptx = output_dir / f"{stem}.pptx"
    return {
        "analyze_json": analyze_json,
        "global_json": global_json,
        "report_json": report_json,
        "out_md": out_md,
        "out_pptx": out_pptx,
    }


def main() -> int:
    try:
        args = parse_args()
        cfg = load_config_json(Path(args.config_json) if args.config_json else None)
        apply_config_defaults(args, cfg)
        validate_args(args)
    except Exception as e:
        print(f"[error] {e}", file=sys.stderr)
        return 2
    global SECTION_PATTERNS
    try:
        SECTION_PATTERNS = load_section_patterns(Path(args.section_map_json) if args.section_map_json else None)
    except Exception as e:
        print(f"[error] invalid --section_map_json: {e}", file=sys.stderr)
        return 2

    def vprint(msg: str) -> None:
        if args.verbose:
            print(msg, flush=True)

    output_dir, file_stem, session_name = resolve_session_context(args)
    args.session_name = session_name
    output_dir.mkdir(parents=True, exist_ok=True)
    status_file = Path(args.status_file) if args.status_file else output_dir / f"{file_stem}.status.json"
    log_file = Path(args.log_file) if args.log_file else output_dir / f"{file_stem}.run.log"
    paper_status_file = (
        Path(args.paper_status_file)
        if args.paper_status_file
        else output_dir / f"{file_stem}.paper_status.jsonl"
    )
    paths = get_stage_paths(args, output_dir, file_stem)
    try:
        paper_status_file.unlink(missing_ok=True)
    except Exception:
        pass
    append_log(
        log_file,
        f"start collection={args.collection} mode={args.mode} llm_mode={args.llm_mode} session={session_name}",
    )
    vprint(f"[start] collection={args.collection} mode={args.mode} llm_mode={args.llm_mode} session={session_name}")
    write_status(
        status_file,
        stage="init",
        status="running",
        progress_current=0,
        progress_total=100,
        message="初始化参数与路径",
        extra={"session_name": session_name, "session_output_dir": str(output_dir)},
    )

    include_images = parse_bool(args.include_images)
    llm_api_key = os.getenv(args.llm_api_key_env, "")
    llm_client = LLMClient(
        mode=args.llm_mode,
        model=args.llm_model,
        base_url=args.llm_base_url,
        api_key=llm_api_key,
        codex_bin=args.codex_bin,
        timeout_sec=args.llm_timeout_sec,
        max_input_chars=args.llm_max_input_chars,
        max_tokens=args.llm_max_tokens,
    )
    if args.llm_mode == "codex_cli" and not llm_client.enabled:
        print(
            "[warn] llm_mode=codex_cli but codex executable not found; fallback to rule-based mode",
            file=sys.stderr,
        )
    if args.llm_mode == "openai_compatible" and not llm_client.enabled:
        print(
            f"[warn] llm mode requested but api key env '{args.llm_api_key_env}' is empty; fallback to rule-based mode",
            file=sys.stderr,
        )
    rag_enabled = parse_bool(args.rag_enabled)
    rag_client = ZoteroMcpRagClient(
        enabled=rag_enabled,
        python_bin=args.rag_python_bin,
        config_path=args.rag_config_path,
        use_local=parse_bool(args.rag_use_local),
        timeout_sec=max(30, int(args.llm_timeout_sec)),
    )
    write_status(
        status_file,
        stage="manifest",
        status="running",
        progress_current=5,
        progress_total=100,
        message="读取collection并构建文献清单",
    )

    analyze_payload: Optional[Dict[str, Any]] = None
    global_payload: Optional[Dict[str, Any]] = None

    if args.mode in {"analyze", "all"}:
        manifest_source = "input_manifest"
        if args.manifest:
            try:
                manifest = read_manifest(Path(args.manifest))
            except Exception as e:
                print(f"[error] invalid manifest: {e}", file=sys.stderr)
                return 2
        else:
            try:
                manifest = build_manifest_from_zotero_collection(
                    collection=args.collection,
                    zotero_db=Path(args.zotero_db),
                    zotero_storage_dir=Path(args.zotero_storage_dir),
                )
                manifest_source = "zotero_db"
                manifest_out = output_dir / f"{file_stem}.manifest.json"
                manifest_out.write_text(json.dumps(manifest, ensure_ascii=False, indent=2), encoding="utf-8")
                print(f"[ok] wrote {manifest_out}")
                append_log(log_file, f"manifest generated: {manifest_out}")
                vprint(f"[manifest] generated: {manifest_out}")
            except Exception as e:
                print(f"[error] failed to build manifest from zotero collection: {e}", file=sys.stderr)
                return 2

        items = [x for x in manifest.get("items", []) if isinstance(x, dict) and x.get("pdf_path")]
        if not items:
            print("[error] no items with pdf_path found", file=sys.stderr)
            return 1

        items = dedup_items(items)[: args.max_papers]
        append_log(log_file, f"papers selected: {len(items)}")
        vprint(f"[papers] selected={len(items)}")
        write_status(
            status_file,
            stage="paper_analysis",
            status="running",
            progress_current=10,
            progress_total=100,
            message=f"开始逐篇解析，共{len(items)}篇",
            extra={"paper_total": len(items)},
        )

        papers = []
        failures = []
        for idx, it in enumerate(items, start=1):
            vprint(f"[paper {idx}/{len(items)}] start: {it.get('title','unknown')[:90]}")
            append_jsonl(
                paper_status_file,
                {
                    "timestamp": ts_now(),
                    "index": idx,
                    "total": len(items),
                    "title": it.get("title", "unknown"),
                    "status": "started",
                    "pdf_path": it.get("pdf_path", ""),
                },
            )
            if not Path(it.get("pdf_path", "")).exists():
                failures.append({"title": it.get("title", "unknown"), "reason": "pdf_not_found", "pdf_path": it.get("pdf_path", "")})
                append_jsonl(
                    paper_status_file,
                    {
                        "timestamp": ts_now(),
                        "index": idx,
                        "total": len(items),
                        "title": it.get("title", "unknown"),
                        "status": "failed",
                        "reason": "pdf_not_found",
                    },
                )
                continue
            p = analyze_paper(
                it,
                output_dir,
                include_images,
                args.language,
                llm_client=llm_client,
                progress_cb=(lambda m, idx=idx, n=len(items): vprint(f"[paper {idx}/{n}] {m}")),
            )
            papers.append(p)
            write_status(
                status_file,
                stage="paper_analysis",
                status="running",
                progress_current=10 + int(65 * idx / max(1, len(items))),
                progress_total=100,
                message=f"解析中 {idx}/{len(items)}: {p.get('title','unknown')[:60]}",
                extra={"paper_current": idx, "paper_total": len(items)},
            )
            if p.get("parse_failed"):
                failures.append({"title": p.get("title", "unknown"), "reason": p.get("failure_reason", "parse_failed")})
            append_jsonl(
                paper_status_file,
                {
                    "timestamp": ts_now(),
                    "index": idx,
                    "total": len(items),
                    "title": p.get("title", "unknown"),
                    "status": "failed" if p.get("parse_failed") else "parsed",
                    "venue": p.get("venue", "unknown"),
                    "year": p.get("year", "unknown"),
                    "authors_n": len(p.get("authors", [])),
                    "has_image": bool(p.get("image_path")),
                    "llm_task_used": bool(p.get("llm_task_used")),
                    "llm_method_used": bool(p.get("llm_method_used")),
                    "llm_contrib_used": bool(p.get("llm_contrib_used")),
                    "llm_limits_used": bool(p.get("llm_limits_used")),
                    "llm_error": p.get("llm_error", ""),
                    "llm_error_debug": p.get("llm_error_debug", ""),
                    "llm_raw_preview": p.get("llm_raw_preview", ""),
                    "keywords": p.get("keywords", []),
                    "task_definition": [x.get("text", "") for x in p.get("task_definition", [])[:1]],
                    "core_method": [x.get("text", "") for x in p.get("core_method", [])[:3]],
                    "main_contributions": [x.get("text", "") for x in p.get("main_contributions", [])[:4]],
                    "limitations": [x.get("text", "") for x in p.get("limitations", [])[:3]],
                },
            )
        analyze_payload = {
            "run_config": vars(args),
            "session_name": session_name,
            "session_output_dir": str(output_dir),
            "collection_info": {"collection": args.collection, "manifest_source": manifest_source},
            "paper_count_total": len(items),
            "paper_count_processed": len([p for p in papers if not p.get("parse_failed")]),
            "papers": papers,
            "failures": failures,
        }
        paths["analyze_json"].write_text(json.dumps(analyze_payload, ensure_ascii=False, indent=2), encoding="utf-8")
        print(f"[ok] wrote {paths['analyze_json']}")
        append_log(log_file, f"analyze written: {paths['analyze_json']}")
        if args.mode == "analyze":
            write_status(
                status_file,
                stage="done",
                status="completed",
                progress_current=100,
                progress_total=100,
                message="analyze完成",
                extra={"analyze_json": str(paths["analyze_json"])},
            )
            return 0

    if args.mode in {"global", "all"}:
        if analyze_payload is None:
            if not paths["analyze_json"].exists():
                print(f"[error] analyze json not found: {paths['analyze_json']}", file=sys.stderr)
                return 2
            analyze_payload = json.loads(paths["analyze_json"].read_text(encoding="utf-8"))
        processable = [p for p in analyze_payload.get("papers", []) if not p.get("parse_failed")]
        if not processable:
            print("[error] no processable papers in analyze json", file=sys.stderr)
            return 1
        write_status(
            status_file,
            stage="synthesis",
            status="running",
            progress_current=82,
            progress_total=100,
            message="聚类与跨论文总结（RAG增强）" if rag_enabled else "聚类与跨论文总结",
        )
        clusters = cluster_papers(processable, forced_k=args.cluster_k)
        overview = make_overview(
            clusters,
            processable,
            args.language,
            llm_client=llm_client,
            rag_client=rag_client if rag_enabled else None,
            rag_top_k=args.rag_top_k,
        )
        final_syn = make_final_synthesis(
            overview,
            args.language,
            llm_client=llm_client,
            rag_enabled=rag_enabled,
        )
        quality = run_quality_gates(processable, overview, final_syn)
        global_payload = {
            "run_config": vars(args),
            "session_name": session_name,
            "session_output_dir": str(output_dir),
            "collection_info": analyze_payload.get("collection_info", {"collection": args.collection}),
            "paper_count_processed": len(processable),
            "cluster_k_requested": args.cluster_k,
            "cluster_k_actual": len(overview),
            "clusters": overview,
            "final_synthesis": final_syn,
            "quality_gate_results": quality,
            "rag": {
                "enabled": rag_enabled,
                "top_k": args.rag_top_k,
                "python_bin": rag_client.python_bin,
                "config_path": args.rag_config_path,
                "use_local": parse_bool(args.rag_use_local),
                "last_error": rag_client.last_error,
            },
        }
        paths["global_json"].write_text(json.dumps(global_payload, ensure_ascii=False, indent=2), encoding="utf-8")
        print(f"[ok] wrote {paths['global_json']}")
        append_log(log_file, f"global written: {paths['global_json']}")
        if args.mode == "global":
            write_status(
                status_file,
                stage="done",
                status="completed",
                progress_current=100,
                progress_total=100,
                message="global完成",
                extra={"global_json": str(paths["global_json"])},
            )
            return 0

    if args.mode in {"render", "all"}:
        if analyze_payload is None:
            if not paths["analyze_json"].exists():
                print(f"[error] analyze json not found: {paths['analyze_json']}", file=sys.stderr)
                return 2
            analyze_payload = json.loads(paths["analyze_json"].read_text(encoding="utf-8"))
        if global_payload is None:
            if not paths["global_json"].exists():
                print(f"[error] global json not found: {paths['global_json']}", file=sys.stderr)
                return 2
            global_payload = json.loads(paths["global_json"].read_text(encoding="utf-8"))
        processable = [p for p in analyze_payload.get("papers", []) if not p.get("parse_failed")]
        overview = global_payload.get("clusters", [])
        final_syn = global_payload.get("final_synthesis", {})

        write_status(
            status_file,
            stage="render",
            status="running",
            progress_current=92,
            progress_total=100,
            message="渲染 Markdown / PPTX",
        )
        render_markdown(paths["out_md"], args.collection, overview, processable, final_syn)
        ppt_warn = render_ppt(paths["out_pptx"], overview, processable, final_syn, args.language)

        deck_outline = {
            "part1_overview_title": "文献整体结构梳理" if args.language == "zh" else "Literature Structure Overview",
            "part2_per_paper_count": len(processable),
            "part3_title": "结构归纳与研究机会" if args.language == "zh" else "Synthesis and Research Opportunities",
        }
        report = {
            "run_config": vars(args),
            "collection_info": analyze_payload.get("collection_info", {"collection": args.collection}),
            "paper_count_total": analyze_payload.get("paper_count_total", len(processable)),
            "paper_count_processed": len(processable),
            "clusters": overview,
            "papers": processable,
            "deck_outline": deck_outline,
            "quality_gate_results": global_payload.get("quality_gate_results", {}),
            "failures": analyze_payload.get("failures", []),
            "analyze_json": str(paths["analyze_json"]),
            "global_json": str(paths["global_json"]),
        }
        paths["report_json"].write_text(json.dumps(report, ensure_ascii=False, indent=2), encoding="utf-8")
        print(f"[ok] wrote {paths['report_json']}")
        print(f"[ok] wrote {paths['out_md']}")
        if ppt_warn:
            print(f"[warn] {ppt_warn}")
        else:
            print(f"[ok] wrote {paths['out_pptx']}")

    append_log(log_file, "completed")
    write_status(
        status_file,
        stage="done",
        status="completed",
        progress_current=100,
        progress_total=100,
        message="任务完成",
        extra={
            "analyze_json": str(paths["analyze_json"]),
            "global_json": str(paths["global_json"]),
            "output_json": str(paths["report_json"]),
            "output_md": str(paths["out_md"]),
            "output_pptx": str(paths["out_pptx"]),
        },
    )
    return 0


if __name__ == "__main__":
    sys.exit(main())
